import os
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from config import UPLOAD_DIR
from models import Portfolio, User

def generate_portfolio_ppt(portfolio: Portfolio, owner: User) -> str:
    """Generates a PowerPoint presentation from a Portfolio and returns its path."""
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = portfolio.name or owner.username
    subtitle.text = f"{portfolio.role or 'Professional Portfolio'}\n\n{portfolio.tagline or ''}"
    
    # Slide 2: About Me (Bio)
    if portfolio.bio:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "About Me"
        tf = body_shape.text_frame
        tf.text = portfolio.bio
        
    # Slide 3: Skills
    if portfolio.skills:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Skills & Expertise"
        tf = body_shape.text_frame
        
        if isinstance(portfolio.skills, list):
            for skill in portfolio.skills:
                p = tf.add_paragraph()
                p.text = str(skill)
        elif isinstance(portfolio.skills, dict):
            for cat, skills in portfolio.skills.items():
                p = tf.add_paragraph()
                p.text = f"{cat}: {', '.join(skills)}"
        
    # Slide 4: Experience
    if portfolio.experience:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Experience"
        tf = body_shape.text_frame
        
        if isinstance(portfolio.experience, list):
            for exp in portfolio.experience:
                title_str = exp.get("title", "")
                company = exp.get("company", "")
                dates = exp.get("dates", "")
                p = tf.add_paragraph()
                p.text = f"{title_str} at {company} ({dates})"
                p.level = 0
                
                desc = exp.get("description", "")
                if desc:
                    p2 = tf.add_paragraph()
                    p2.text = desc
                    p2.level = 1

    # Slide 5: Projects
    if portfolio.projects:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Projects"
        tf = body_shape.text_frame
        
        if isinstance(portfolio.projects, list):
            for proj in portfolio.projects:
                p = tf.add_paragraph()
                p.text = proj.get("name", "Project")
                p.level = 0
                
                desc = proj.get("description", "")
                if desc:
                    p2 = tf.add_paragraph()
                    p2.text = desc
                    p2.level = 1
                    
    # Slide 6: Contact
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Contact Information"
    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = f"Email: {owner.email}"
    
    if portfolio.contact and isinstance(portfolio.contact, dict):
        for key, val in portfolio.contact.items():
            if val:
                p = tf.add_paragraph()
                p.text = f"{key.title()}: {val}"

    # Save presentation
    filename = f"{owner.username}_portfolio_{portfolio.id}.pptx"
    filepath = os.path.join(UPLOAD_DIR, filename)
    prs.save(filepath)
    
    return filepath

def generate_ats_resume_docx_from_data(ats_data: dict, slug: str) -> str:
    """Generates an ATS-friendly minimal Word document resume using extracted structured data."""
    document = Document()
    
    # Set document margins (1 inch is standard for ATS)
    for section in document.sections:
        section.top_margin = DocxInches(1.0)
        section.bottom_margin = DocxInches(1.0)
        section.left_margin = DocxInches(1.0)
        section.right_margin = DocxInches(1.0)
        
    def add_heading(text):
        p = document.add_heading(text, level=1)
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
        
    # --- Header (Contact info) ---
    personal = ats_data.get("personal_info", {})
    name = personal.get("name", "Name")
    
    name_p = document.add_paragraph()
    name_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    name_run = name_p.add_run(name)
    name_run.bold = True
    name_run.font.size = DocxPt(16)
    
    contacts = []
    if personal.get("email"): contacts.append(personal.get("email"))
    if personal.get("phone"): contacts.append(personal.get("phone"))
    if personal.get("linkedin"): contacts.append(personal.get("linkedin"))
    if personal.get("github"): contacts.append(personal.get("github"))
    if personal.get("portfolio"): contacts.append(personal.get("portfolio"))
    
    if contacts:
        contact_p = document.add_paragraph(" | ".join(contacts))
        contact_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
    document.add_paragraph() # Spacing
    
    # --- Summary ---
    summary = ats_data.get("professional_summary", "")
    if summary:
        add_heading("Professional Summary")
        document.add_paragraph(summary)
        
    # --- Skills ---
    skills = ats_data.get("skills", {})
    if skills:
        add_heading("Skills")
        if isinstance(skills, dict):
            for cat, items in skills.items():
                p = document.add_paragraph()
                p.add_run(f"{cat}: ").bold = True
                p.add_run(", ".join(items))
        elif isinstance(skills, list):
            document.add_paragraph(", ".join(skills))
            
    # --- Experience ---
    experience = ats_data.get("experience", [])
    if experience:
        add_heading("Experience")
        for exp in experience:
            p = document.add_paragraph()
            p.add_run(exp.get("role", "Role")).bold = True
            if exp.get('company'):
                p.add_run(f" | {exp.get('company')}")
            if exp.get('location'):
                p.add_run(f", {exp.get('location')}")
            if exp.get('date'):
                p.add_run(f" | {exp.get('date')}")
                
            for bullet in exp.get("bullets", []):
                document.add_paragraph(bullet, style='List Bullet')
                
    # --- Projects ---
    projects = ats_data.get("projects", [])
    if projects:
        add_heading("Projects")
        for proj in projects:
            p = document.add_paragraph()
            p.add_run(proj.get("name", "Project")).bold = True
            tech = proj.get("technologies")
            if tech:
                if isinstance(tech, list):
                    tech_str = ", ".join(tech)
                else:
                    tech_str = str(tech)
                p.add_run(f" ({tech_str})")
            
            desc = proj.get("description", "")
            if desc:
                document.add_paragraph(desc, style='List Bullet')
                
            for bullet in proj.get("bullets", []):
                document.add_paragraph(bullet, style='List Bullet')

    # --- Education ---
    education = ats_data.get("education", [])
    if education:
        add_heading("Education")
        for edu in education:
            p = document.add_paragraph()
            p.add_run(edu.get("degree", "")).bold = True
            if edu.get('institution'):
                p.add_run(f" | {edu.get('institution')}")
            if edu.get('date'):
                p.add_run(f" | {edu.get('date')}")
            if edu.get('details'):
                document.add_paragraph(edu.get('details'))

    # --- Certifications ---
    certs = ats_data.get("certifications", [])
    if certs:
        add_heading("Certifications")
        for cert in certs:
            document.add_paragraph(cert, style='List Bullet')
                
    # Save document
    filename = f"{slug}_ATS_Resume.docx"
    filepath = os.path.join(UPLOAD_DIR, filename)
    document.save(filepath)
    
    return filepath
